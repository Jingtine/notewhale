from datetime import datetime, timedelta
from pathlib import Path
from typing import Optional
import base64
import hashlib
import hmac
import json
import mimetypes
import os
import secrets
import urllib.request

from fastapi import Depends, FastAPI, File, Header, HTTPException, Query, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from sqlalchemy import inspect, text, or_
from sqlalchemy.orm import Session

from database import Base, engine, get_db
from models import Course, DDL, Folder, Note, Resource, User

BASE_DIR = Path(__file__).resolve().parent
UPLOAD_DIR = BASE_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

Base.metadata.create_all(bind=engine)

SECRET_KEY = os.getenv("NOTEWHALE_SECRET_KEY", "notewhale-local-dev-secret-change-before-deploy")
TOKEN_EXPIRE_DAYS = int(os.getenv("NOTEWHALE_TOKEN_EXPIRE_DAYS", "14"))

VISION_API_URL = os.getenv("NOTEWHALE_VISION_API_URL", "").strip()
VISION_API_KEY = os.getenv("NOTEWHALE_VISION_API_KEY", "").strip()
VISION_MODEL = os.getenv("NOTEWHALE_VISION_MODEL", "").strip()

TEXT_API_URL = os.getenv("NOTEWHALE_TEXT_API_URL", VISION_API_URL).strip()
TEXT_API_KEY = os.getenv("NOTEWHALE_TEXT_API_KEY", VISION_API_KEY).strip()
TEXT_MODEL = os.getenv("NOTEWHALE_TEXT_MODEL", "glm-4-flash-250414").strip()


def ensure_schema():
    """兼容旧版 SQLite 表结构：create_all 不会自动给旧表补字段。"""
    inspector = inspect(engine)

    with engine.begin() as conn:
        table_names = inspector.get_table_names()

        if "users" not in table_names:
            Base.metadata.create_all(bind=engine)
            inspector = inspect(engine)
            table_names = inspector.get_table_names()

        if "folders" in table_names:
            folder_columns = {column["name"] for column in inspector.get_columns("folders")}
            if "user_id" not in folder_columns:
                conn.execute(text("ALTER TABLE folders ADD COLUMN user_id INTEGER"))

        if "courses" in table_names:
            course_columns = {column["name"] for column in inspector.get_columns("courses")}

            if "folder_id" not in course_columns:
                conn.execute(text("ALTER TABLE courses ADD COLUMN folder_id INTEGER"))

            if "folder_name" not in course_columns:
                conn.execute(text("ALTER TABLE courses ADD COLUMN folder_name VARCHAR DEFAULT ''"))

            if "is_deleted" not in course_columns:
                conn.execute(text("ALTER TABLE courses ADD COLUMN is_deleted BOOLEAN DEFAULT 0"))

            if "deleted_at" not in course_columns:
                conn.execute(text("ALTER TABLE courses ADD COLUMN deleted_at DATETIME"))

            if "deleted_folder_id" not in course_columns:
                conn.execute(text("ALTER TABLE courses ADD COLUMN deleted_folder_id INTEGER"))

            if "deleted_folder_title" not in course_columns:
                conn.execute(text("ALTER TABLE courses ADD COLUMN deleted_folder_title VARCHAR DEFAULT ''"))

            if "user_id" not in course_columns:
                conn.execute(text("ALTER TABLE courses ADD COLUMN user_id INTEGER"))

        if "ddls" in table_names:
            ddl_columns = {column["name"] for column in inspector.get_columns("ddls")}
            if "user_id" not in ddl_columns:
                conn.execute(text("ALTER TABLE ddls ADD COLUMN user_id INTEGER"))

        if "notes" in table_names:
            note_columns = {column["name"] for column in inspector.get_columns("notes")}
            if "user_id" not in note_columns:
                conn.execute(text("ALTER TABLE notes ADD COLUMN user_id INTEGER"))

        if "resources" in table_names:
            resource_columns = {column["name"] for column in inspector.get_columns("resources")}
            if "user_id" not in resource_columns:
                conn.execute(text("ALTER TABLE resources ADD COLUMN user_id INTEGER"))


ensure_schema()

app = FastAPI(
    title="NoteWhale API",
    description="鲸记 NoteWhale 后端接口",
    version="0.6.0-ai",
)

# Frontend origins allowed to call this API.
# Local development:
#   http://localhost:5173
#   http://127.0.0.1:5173
# Production:
#   https://notewhale.vercel.app
# You can also add more origins through NOTEWHALE_FRONTEND_ORIGINS,
# separated by commas.
frontend_origins = [
    "http://localhost:5173",
    "http://127.0.0.1:5173",
    "https://notewhale.vercel.app",
]

extra_origins = os.getenv("NOTEWHALE_FRONTEND_ORIGINS", "").strip()
if extra_origins:
    frontend_origins.extend(
        origin.strip().rstrip("/")
        for origin in extra_origins.split(",")
        if origin.strip()
    )

app.add_middleware(
    CORSMiddleware,
    allow_origins=frontend_origins,
    allow_origin_regex=r"^https://notewhale[-a-zA-Z0-9]*-jingtines-projects\.vercel\.app$",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.mount("/uploads", StaticFiles(directory=str(UPLOAD_DIR)), name="uploads")


class RegisterRequest(BaseModel):
    account: str
    password: str
    name: str


class LoginRequest(BaseModel):
    account: str
    password: str


class FolderCreate(BaseModel):
    title: str


class FolderUpdate(BaseModel):
    title: Optional[str] = None


class CourseCreate(BaseModel):
    title: str
    starred: bool = False
    folderId: Optional[int] = None
    folderName: str = ""


class CourseUpdate(BaseModel):
    title: Optional[str] = None
    starred: Optional[bool] = None
    folderId: Optional[int] = None
    folderName: Optional[str] = None


class RestoreCourseRequest(BaseModel):
    folderId: Optional[int] = None
    folderName: str = ""


class DDLCreate(BaseModel):
    title: str
    date: str
    courseId: Optional[int] = None
    courseName: str = "未归属课程"
    platform: str = ""
    note: str = ""
    completed: bool = False
    source: str = "手动新建"


class DDLUpdate(BaseModel):
    title: Optional[str] = None
    date: Optional[str] = None
    courseId: Optional[int] = None
    courseName: Optional[str] = None
    platform: Optional[str] = None
    note: Optional[str] = None
    completed: Optional[bool] = None
    source: Optional[str] = None


class NoteCreate(BaseModel):
    title: str
    content: str = ""
    courseId: Optional[int] = None
    courseName: str = ""
    source: str = "手动记录"
    aiGenerated: bool = False


class NoteUpdate(BaseModel):
    title: Optional[str] = None
    content: Optional[str] = None
    courseId: Optional[int] = None
    courseName: Optional[str] = None
    source: Optional[str] = None
    aiGenerated: Optional[bool] = None


class GenerateNoteRequest(BaseModel):
    courseId: Optional[int] = None
    courseName: str = "课程"
    resourceName: str = "课程资料"
    resourceId: Optional[int] = None
    rawText: str = ""
    noteStyle: str = "复习型"


class RecognizeDDLRequest(BaseModel):
    courseId: Optional[int] = None
    courseName: str = "未归属课程"


def timestamp(dt):
    return int(dt.timestamp() * 1000) if dt else None


def normalize_account(account: str):
    return (account or "").strip().lower()


def hash_password(password: str):
    if not password or len(password) < 1:
        raise HTTPException(status_code=400, detail="密码不能为空")

    salt = secrets.token_hex(16)
    iterations = 120000
    digest = hashlib.pbkdf2_hmac(
        "sha256",
        password.encode("utf-8"),
        salt.encode("utf-8"),
        iterations,
    ).hex()

    return f"pbkdf2_sha256${iterations}${salt}${digest}"


def verify_password(password: str, stored_hash: str):
    try:
        algorithm, iterations_text, salt, expected_digest = stored_hash.split("$", 3)
        if algorithm != "pbkdf2_sha256":
            return False

        iterations = int(iterations_text)
        digest = hashlib.pbkdf2_hmac(
            "sha256",
            password.encode("utf-8"),
            salt.encode("utf-8"),
            iterations,
        ).hex()

        return hmac.compare_digest(digest, expected_digest)
    except Exception:
        return False


def b64url_encode(data: bytes):
    return base64.urlsafe_b64encode(data).decode("utf-8").rstrip("=")


def b64url_decode(data: str):
    padding = "=" * (-len(data) % 4)
    return base64.urlsafe_b64decode((data + padding).encode("utf-8"))


def create_access_token(user: User):
    expire_at = datetime.utcnow() + timedelta(days=TOKEN_EXPIRE_DAYS)
    payload = {
        "sub": user.id,
        "account": user.account,
        "exp": int(expire_at.timestamp()),
    }

    payload_part = b64url_encode(json.dumps(payload, separators=(",", ":")).encode("utf-8"))
    signature = hmac.new(
        SECRET_KEY.encode("utf-8"),
        payload_part.encode("utf-8"),
        hashlib.sha256,
    ).digest()

    return f"{payload_part}.{b64url_encode(signature)}"


def parse_access_token(token: str):
    try:
        payload_part, signature_part = token.split(".", 1)

        expected_signature = hmac.new(
            SECRET_KEY.encode("utf-8"),
            payload_part.encode("utf-8"),
            hashlib.sha256,
        ).digest()

        if not hmac.compare_digest(b64url_encode(expected_signature), signature_part):
            raise ValueError("bad signature")

        payload = json.loads(b64url_decode(payload_part).decode("utf-8"))

        if int(payload.get("exp", 0)) < int(datetime.utcnow().timestamp()):
            raise ValueError("expired")

        return payload
    except Exception:
        raise HTTPException(status_code=401, detail="登录状态已失效，请重新登录")


def user_to_dict(user: User):
    return {
        "id": user.id,
        "account": user.account,
        "email": user.account if "@" in user.account else "",
        "name": user.name,
        "avatar": user.avatar or (user.name or "鲸")[:1],
        "createdAt": timestamp(user.created_at),
        "authMode": "api",
    }


def auth_response(user: User):
    return {
        "token": create_access_token(user),
        "user": user_to_dict(user),
    }


def get_current_user(
    authorization: Optional[str] = Header(None),
    db: Session = Depends(get_db),
):
    if not authorization or not authorization.lower().startswith("bearer "):
        raise HTTPException(status_code=401, detail="请先登录账号")

    token = authorization.split(" ", 1)[1].strip()
    payload = parse_access_token(token)

    user = db.query(User).filter(User.id == payload.get("sub")).first()
    if not user:
        raise HTTPException(status_code=401, detail="账号不存在，请重新登录")

    return user


def folder_to_dict(folder: Folder, courses=None):
    return {
        "id": folder.id,
        "title": folder.title,
        "createdAt": timestamp(folder.created_at),
        "courses": courses if courses is not None else [],
    }


def course_to_dict(course: Course):
    return {
        "id": course.id,
        "title": course.title,
        "starred": course.starred,
        "folderId": course.folder_id,
        "folderName": course.folder_name or "",
        "isDeleted": bool(getattr(course, "is_deleted", False)),
        "deletedAt": timestamp(getattr(course, "deleted_at", None)),
        "deletedFolderId": getattr(course, "deleted_folder_id", None),
        "deletedFolderTitle": getattr(course, "deleted_folder_title", "") or "",
        "createdAt": timestamp(course.created_at),
    }


def ddl_to_dict(ddl: DDL):
    return {
        "id": ddl.id,
        "title": ddl.title,
        "date": ddl.date,
        "courseId": ddl.course_id,
        "courseName": ddl.course_name,
        "platform": ddl.platform,
        "note": ddl.note,
        "completed": ddl.completed,
        "source": ddl.source,
        "createdAt": timestamp(ddl.created_at),
    }


def note_to_dict(note: Note):
    return {
        "id": note.id,
        "title": note.title,
        "content": note.content,
        "courseId": note.course_id,
        "courseName": note.course_name,
        "source": note.source,
        "aiGenerated": note.ai_generated,
        "createdAt": timestamp(note.created_at),
        "updatedAt": timestamp(note.updated_at),
    }


def resource_to_dict(resource: Resource):
    filename = Path(resource.file_path).name
    return {
        "id": resource.id,
        "name": resource.name,
        "type": resource.file_type,
        "filePath": resource.file_path,
        "url": f"/uploads/{filename}",
        "size": resource.size,
        "courseId": resource.course_id,
        "courseName": resource.course_name,
        "createdAt": timestamp(resource.created_at),
    }


def get_file_type(filename: str):
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return "PDF"
    if ext in {"ppt", "pptx"}:
        return "PPT"
    if ext in {"doc", "docx"}:
        return "Word"
    if ext in {"jpg", "jpeg", "png", "webp"}:
        return "图片"
    if ext in {"mp3", "wav", "m4a"}:
        return "录音"
    if ext in {"txt", "md"}:
        return "文本"
    return "资料"


def active_course_filter():
    return or_(Course.is_deleted == False, Course.is_deleted.is_(None))


def soft_delete_course(course: Course, folder_id=None, folder_title: str = ""):
    course.is_deleted = True
    course.deleted_at = datetime.utcnow()
    course.deleted_folder_id = folder_id if folder_id is not None else course.folder_id
    course.deleted_folder_title = folder_title or course.folder_name or ""


def ensure_folder(db: Session, folder_id: Optional[int], folder_name: str, user: User):
    if folder_id:
        folder = (
            db.query(Folder)
            .filter(Folder.id == folder_id, Folder.user_id == user.id)
            .first()
        )
        if folder:
            return folder
        raise HTTPException(status_code=404, detail="目标文件夹不存在或不属于当前账号")

    title = (folder_name or "").strip()
    if not title:
        return None

    folder = (
        db.query(Folder)
        .filter(Folder.title == title, Folder.user_id == user.id)
        .first()
    )
    if folder:
        return folder

    folder = Folder(title=title, user_id=user.id)
    db.add(folder)
    db.commit()
    db.refresh(folder)
    return folder


def ensure_course_access(db: Session, course_id: Optional[int], user: User):
    if not course_id:
        return None

    course = (
        db.query(Course)
        .filter(Course.id == course_id, Course.user_id == user.id, active_course_filter())
        .first()
    )
    if not course:
        raise HTTPException(status_code=404, detail="目标课程不存在或不属于当前账号")

    return course


@app.get("/health")
def health():
    return {
        "status": "ok",
        "service": "notewhale-backend",
        "version": "0.6.0-ai",
    }


@app.post("/api/auth/register")
def register(payload: RegisterRequest, db: Session = Depends(get_db)):
    account = normalize_account(payload.account)
    name = (payload.name or "").strip()
    password = payload.password or ""

    if not account:
        raise HTTPException(status_code=400, detail="账号不能为空")
    if not name:
        raise HTTPException(status_code=400, detail="昵称不能为空")
    if len(password) < 6:
        raise HTTPException(status_code=400, detail="密码至少需要 6 位")

    exists = db.query(User).filter(User.account == account).first()
    if exists:
        raise HTTPException(status_code=409, detail="该账号已注册，请直接登录")

    user = User(
        account=account,
        name=name,
        password_hash=hash_password(password),
        avatar=name[:1].upper(),
    )

    db.add(user)
    db.commit()
    db.refresh(user)

    return auth_response(user)


@app.post("/api/auth/login")
def login(payload: LoginRequest, db: Session = Depends(get_db)):
    account = normalize_account(payload.account)
    password = payload.password or ""

    user = db.query(User).filter(User.account == account).first()
    if not user or not verify_password(password, user.password_hash):
        raise HTTPException(status_code=401, detail="账号或密码错误")

    return auth_response(user)


@app.get("/api/auth/me")
def me(current_user: User = Depends(get_current_user)):
    return user_to_dict(current_user)


@app.get("/api/folders")
def list_folders(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    folders = (
        db.query(Folder)
        .filter(Folder.user_id == current_user.id)
        .order_by(Folder.created_at.asc())
        .all()
    )
    result = []

    for folder in folders:
        courses = (
            db.query(Course)
            .filter(
                Course.user_id == current_user.id,
                Course.folder_id == folder.id,
                active_course_filter(),
            )
            .order_by(Course.created_at.desc())
            .all()
        )
        result.append(folder_to_dict(folder, [course_to_dict(course) for course in courses]))

    unassigned_courses = (
        db.query(Course)
        .filter(
            Course.user_id == current_user.id,
            Course.folder_id.is_(None),
            active_course_filter(),
        )
        .order_by(Course.created_at.desc())
        .all()
    )

    if unassigned_courses:
        result.append(
            {
                "id": None,
                "title": "未归属课程",
                "createdAt": None,
                "courses": [course_to_dict(course) for course in unassigned_courses],
            }
        )

    return result


@app.post("/api/folders")
def create_folder(
    folder: FolderCreate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    title = folder.title.strip()
    if not title:
        raise HTTPException(status_code=400, detail="文件夹名称不能为空")

    new_folder = Folder(title=title, user_id=current_user.id)
    db.add(new_folder)
    db.commit()
    db.refresh(new_folder)
    return folder_to_dict(new_folder)


@app.put("/api/folders/{folder_id}")
def update_folder(
    folder_id: int,
    payload: FolderUpdate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    folder = (
        db.query(Folder)
        .filter(Folder.id == folder_id, Folder.user_id == current_user.id)
        .first()
    )
    if not folder:
        raise HTTPException(status_code=404, detail="文件夹不存在")

    if payload.title is not None:
        title = payload.title.strip()
        if not title:
            raise HTTPException(status_code=400, detail="文件夹名称不能为空")

        folder.title = title
        db.query(Course).filter(
            Course.folder_id == folder.id,
            Course.user_id == current_user.id,
        ).update({Course.folder_name: title}, synchronize_session=False)

    db.commit()
    db.refresh(folder)
    return folder_to_dict(folder)


@app.delete("/api/folders/{folder_id}")
def delete_folder(
    folder_id: int,
    deleteCourses: bool = Query(False),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    folder = (
        db.query(Folder)
        .filter(Folder.id == folder_id, Folder.user_id == current_user.id)
        .first()
    )
    if not folder:
        raise HTTPException(status_code=404, detail="文件夹不存在")

    courses = (
        db.query(Course)
        .filter(
            Course.user_id == current_user.id,
            Course.folder_id == folder.id,
            active_course_filter(),
        )
        .all()
    )

    if deleteCourses:
        for course in courses:
            soft_delete_course(course, folder.id, folder.title)
    else:
        for course in courses:
            course.folder_id = None
            course.folder_name = ""

    db.delete(folder)
    db.commit()
    return {"ok": True, "deletedFolderId": folder_id, "deleteCourses": deleteCourses}


@app.post("/api/folders/{folder_id}/delete")
def delete_folder_via_post(
    folder_id: int,
    deleteCourses: bool = Query(True),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    folder = (
        db.query(Folder)
        .filter(Folder.id == folder_id, Folder.user_id == current_user.id)
        .first()
    )
    if not folder:
        raise HTTPException(status_code=404, detail="文件夹不存在")

    courses = (
        db.query(Course)
        .filter(
            Course.user_id == current_user.id,
            Course.folder_id == folder.id,
            active_course_filter(),
        )
        .all()
    )

    if deleteCourses:
        for course in courses:
            soft_delete_course(course, folder.id, folder.title)
    else:
        for course in courses:
            course.folder_id = None
            course.folder_name = ""

    db.delete(folder)
    db.commit()
    return {"ok": True, "deletedFolderId": folder_id, "deleteCourses": deleteCourses}


@app.get("/api/courses")
def list_courses(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    courses = (
        db.query(Course)
        .filter(Course.user_id == current_user.id, active_course_filter())
        .order_by(Course.created_at.desc())
        .all()
    )
    return [course_to_dict(course) for course in courses]


@app.post("/api/courses")
def create_course(
    course: CourseCreate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    title = course.title.strip()
    if not title:
        raise HTTPException(status_code=400, detail="课程名称不能为空")

    folder = ensure_folder(db, course.folderId, course.folderName, current_user)

    new_course = Course(
        title=title,
        starred=course.starred,
        folder_id=folder.id if folder else None,
        folder_name=folder.title if folder else (course.folderName or ""),
        user_id=current_user.id,
    )

    db.add(new_course)
    db.commit()
    db.refresh(new_course)
    return course_to_dict(new_course)


@app.put("/api/courses/{course_id}")
def update_course(
    course_id: int,
    payload: CourseUpdate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    course = (
        db.query(Course)
        .filter(Course.id == course_id, Course.user_id == current_user.id)
        .first()
    )
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    if payload.title is not None:
        title = payload.title.strip()
        if not title:
            raise HTTPException(status_code=400, detail="课程名称不能为空")
        course.title = title

        db.query(DDL).filter(DDL.course_id == course.id, DDL.user_id == current_user.id).update(
            {DDL.course_name: title},
            synchronize_session=False,
        )
        db.query(Note).filter(Note.course_id == course.id, Note.user_id == current_user.id).update(
            {Note.course_name: title},
            synchronize_session=False,
        )
        db.query(Resource).filter(Resource.course_id == course.id, Resource.user_id == current_user.id).update(
            {Resource.course_name: title},
            synchronize_session=False,
        )

    if payload.starred is not None:
        course.starred = payload.starred

    if payload.folderId is not None or payload.folderName is not None:
        folder = ensure_folder(db, payload.folderId, payload.folderName or "", current_user)
        course.folder_id = folder.id if folder else None
        course.folder_name = folder.title if folder else (payload.folderName or "")

    db.commit()
    db.refresh(course)
    return course_to_dict(course)


@app.delete("/api/courses/{course_id}")
def delete_course(
    course_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    course = (
        db.query(Course)
        .filter(Course.id == course_id, Course.user_id == current_user.id)
        .first()
    )
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    if course.is_deleted:
        return {"ok": True, "deletedCourseId": course_id, "softDeleted": True}

    folder_title = course.folder_name or ""
    if course.folder_id:
        folder = (
            db.query(Folder)
            .filter(Folder.id == course.folder_id, Folder.user_id == current_user.id)
            .first()
        )
        if folder:
            folder_title = folder.title

    soft_delete_course(course, course.folder_id, folder_title)
    db.commit()
    db.refresh(course)

    return {"ok": True, "deletedCourseId": course_id, "softDeleted": True, "course": course_to_dict(course)}


@app.get("/api/trash/courses")
def list_deleted_courses(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    courses = (
        db.query(Course)
        .filter(Course.user_id == current_user.id, Course.is_deleted == True)
        .order_by(Course.deleted_at.desc())
        .all()
    )
    return [course_to_dict(course) for course in courses]


@app.post("/api/trash/courses/{course_id}/restore")
def restore_deleted_course(
    course_id: int,
    payload: RestoreCourseRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    course = (
        db.query(Course)
        .filter(Course.id == course_id, Course.user_id == current_user.id)
        .first()
    )
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    folder_name = (payload.folderName or course.deleted_folder_title or course.folder_name or "").strip()
    folder = ensure_folder(db, payload.folderId, folder_name, current_user)

    course.is_deleted = False
    course.deleted_at = None
    course.deleted_folder_id = None
    course.deleted_folder_title = ""
    course.folder_id = folder.id if folder else None
    course.folder_name = folder.title if folder else ""

    db.query(DDL).filter(DDL.course_id == course.id, DDL.user_id == current_user.id).update(
        {DDL.course_name: course.title},
        synchronize_session=False,
    )
    db.query(Note).filter(Note.course_id == course.id, Note.user_id == current_user.id).update(
        {Note.course_name: course.title},
        synchronize_session=False,
    )
    db.query(Resource).filter(Resource.course_id == course.id, Resource.user_id == current_user.id).update(
        {Resource.course_name: course.title},
        synchronize_session=False,
    )

    db.commit()
    db.refresh(course)
    return course_to_dict(course)


@app.delete("/api/trash/courses/{course_id}/permanent")
def permanently_delete_course(
    course_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    course = (
        db.query(Course)
        .filter(Course.id == course_id, Course.user_id == current_user.id)
        .first()
    )
    if not course:
        raise HTTPException(status_code=404, detail="课程不存在")

    db.query(DDL).filter(DDL.course_id == course.id, DDL.user_id == current_user.id).update(
        {DDL.course_id: None, DDL.course_name: "未归属课程"},
        synchronize_session=False,
    )
    db.query(Note).filter(Note.course_id == course.id, Note.user_id == current_user.id).update(
        {Note.course_id: None, Note.course_name: "未归属课程"},
        synchronize_session=False,
    )
    db.query(Resource).filter(Resource.course_id == course.id, Resource.user_id == current_user.id).update(
        {Resource.course_id: None, Resource.course_name: "未归属课程"},
        synchronize_session=False,
    )

    db.delete(course)
    db.commit()
    return {"ok": True, "permanentDeletedCourseId": course_id}


@app.get("/api/ddls")
def list_ddls(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    ddls = (
        db.query(DDL)
        .filter(DDL.user_id == current_user.id)
        .order_by(DDL.created_at.desc())
        .all()
    )
    return [ddl_to_dict(ddl) for ddl in ddls]


@app.post("/api/ddls")
def create_ddl(
    ddl: DDLCreate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    ensure_course_access(db, ddl.courseId, current_user)

    new_ddl = DDL(
        title=ddl.title,
        date=ddl.date,
        course_id=ddl.courseId,
        course_name=ddl.courseName,
        platform=ddl.platform,
        note=ddl.note,
        completed=ddl.completed,
        source=ddl.source,
        user_id=current_user.id,
    )

    db.add(new_ddl)
    db.commit()
    db.refresh(new_ddl)
    return ddl_to_dict(new_ddl)


@app.put("/api/ddls/{ddl_id}")
def update_ddl(
    ddl_id: int,
    payload: DDLUpdate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    ddl = (
        db.query(DDL)
        .filter(DDL.id == ddl_id, DDL.user_id == current_user.id)
        .first()
    )
    if not ddl:
        raise HTTPException(status_code=404, detail="DDL 不存在")

    if payload.courseId is not None:
        ensure_course_access(db, payload.courseId, current_user)

    if payload.title is not None:
        ddl.title = payload.title
    if payload.date is not None:
        ddl.date = payload.date
    if payload.courseId is not None:
        ddl.course_id = payload.courseId
    if payload.courseName is not None:
        ddl.course_name = payload.courseName
    if payload.platform is not None:
        ddl.platform = payload.platform
    if payload.note is not None:
        ddl.note = payload.note
    if payload.completed is not None:
        ddl.completed = payload.completed
    if payload.source is not None:
        ddl.source = payload.source

    db.commit()
    db.refresh(ddl)
    return ddl_to_dict(ddl)


@app.delete("/api/ddls/{ddl_id}")
def delete_ddl(
    ddl_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    ddl = (
        db.query(DDL)
        .filter(DDL.id == ddl_id, DDL.user_id == current_user.id)
        .first()
    )
    if not ddl:
        raise HTTPException(status_code=404, detail="DDL 不存在")

    db.delete(ddl)
    db.commit()
    return {"ok": True, "deletedDdlId": ddl_id}


@app.get("/api/notes")
def list_notes(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    notes = (
        db.query(Note)
        .filter(Note.user_id == current_user.id)
        .order_by(Note.updated_at.desc())
        .all()
    )
    return [note_to_dict(note) for note in notes]


@app.post("/api/notes")
def create_note(
    note: NoteCreate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    ensure_course_access(db, note.courseId, current_user)

    new_note = Note(
        title=note.title,
        content=note.content,
        course_id=note.courseId,
        course_name=note.courseName,
        source=note.source,
        ai_generated=note.aiGenerated,
        user_id=current_user.id,
    )

    db.add(new_note)
    db.commit()
    db.refresh(new_note)
    return note_to_dict(new_note)


@app.put("/api/notes/{note_id}")
def update_note(
    note_id: int,
    payload: NoteUpdate,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = (
        db.query(Note)
        .filter(Note.id == note_id, Note.user_id == current_user.id)
        .first()
    )
    if not note:
        raise HTTPException(status_code=404, detail="笔记不存在")

    if payload.courseId is not None:
        ensure_course_access(db, payload.courseId, current_user)

    if payload.title is not None:
        note.title = payload.title
    if payload.content is not None:
        note.content = payload.content
    if payload.courseId is not None:
        note.course_id = payload.courseId
    if payload.courseName is not None:
        note.course_name = payload.courseName
    if payload.source is not None:
        note.source = payload.source
    if payload.aiGenerated is not None:
        note.ai_generated = payload.aiGenerated

    note.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(note)
    return note_to_dict(note)


@app.delete("/api/notes/{note_id}")
def delete_note(
    note_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    note = (
        db.query(Note)
        .filter(Note.id == note_id, Note.user_id == current_user.id)
        .first()
    )
    if not note:
        raise HTTPException(status_code=404, detail="笔记不存在")

    db.delete(note)
    db.commit()
    return {"ok": True, "deletedNoteId": note_id}



def read_plain_text_file(file_path: Path):
    try:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        try:
            return file_path.read_text(encoding="gbk", errors="ignore")
        except Exception:
            return ""


def normalize_special_symbols(value: str):
    """统一常见公式 / 逻辑 / 数学符号，减少生成笔记时的乱码和 LaTeX 残留。"""
    if not value:
        return ""

    text_value = str(value)

    replacements = {
        "\\leq": "≤",
        "\\le": "≤",
        "<=": "≤",
        "≤": "≤",
        "\\geq": "≥",
        "\\ge": "≥",
        ">=": "≥",
        "≥": "≥",
        "\\neq": "≠",
        "!=": "≠",
        "\\approx": "≈",
        "\\equiv": "≡",
        "\\pm": "±",
        "\\times": "×",
        "\\div": "÷",
        "\\cdot": "·",
        "\\rightarrow": "→",
        "\\Rightarrow": "⇒",
        "\\leftarrow": "←",
        "\\leftrightarrow": "↔",
        "\\Leftrightarrow": "⇔",
        "\\to": "→",
        "\\infty": "∞",
        "\\sum": "∑",
        "\\int": "∫",
        "\\sqrt": "√",
        "\\in": "∈",
        "\\notin": "∉",
        "\\subseteq": "⊆",
        "\\subset": "⊂",
        "\\cup": "∪",
        "\\cap": "∩",
        "\\emptyset": "∅",
        "\\forall": "∀",
        "\\exists": "∃",
        "\\therefore": "∴",
        "\\because": "∵",
        "\\partial": "∂",
        "\\nabla": "∇",
        "\\Delta": "Δ",
        "\\alpha": "α",
        "\\beta": "β",
        "\\gamma": "γ",
        "\\lambda": "λ",
        "\\mu": "μ",
        "\\pi": "π",
        "\\theta": "θ",
        "\\omega": "ω",
        "∵": "∵",
        "∴": "∴",
    }

    for old, new in replacements.items():
        text_value = text_value.replace(old, new)

    # 常见全角符号归一
    text_value = (
        text_value.replace("－", "-")
        .replace("—", "—")
        .replace("×", "×")
        .replace("÷", "÷")
        .replace("（", "(")
        .replace("）", ")")
        .replace("，", "，")
        .replace("；", "；")
        .replace("：", "：")
    )

    return text_value


def extract_pptx_shape_text(shape):
    """递归提取 PPTX 文本框、表格、组合形状中的文字。"""
    parts = []

    try:
        if hasattr(shape, "text"):
            value = normalize_special_symbols((shape.text or "").strip())
            if value:
                parts.append(value)
    except Exception:
        pass

    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = []
                for cell in row.cells:
                    value = normalize_special_symbols((cell.text or "").strip())
                    if value:
                        cells.append(value)
                if cells:
                    parts.append(" | ".join(cells))
    except Exception:
        pass

    try:
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                child_text = extract_pptx_shape_text(child)
                if child_text:
                    parts.append(child_text)
    except Exception:
        pass

    return "\n".join(parts)


def read_pdf_text(file_path: Path):
    """读取 PDF 文本。需要安装 pypdf。"""
    try:
        from pypdf import PdfReader
    except Exception:
        return ""

    try:
        reader = PdfReader(str(file_path))
        parts = []

        for index, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                parts.append(f"[第 {index} 页]\n{page_text}")

        return "\n\n".join(parts)
    except Exception:
        return ""


def read_pptx_text(file_path: Path):
    """读取 PPTX 文本。支持文本框、表格、组合形状、备注页。需要安装 python-pptx。"""
    try:
        from pptx import Presentation
    except Exception:
        return ""

    try:
        prs = Presentation(str(file_path))
        slides = []

        for index, slide in enumerate(prs.slides, start=1):
            texts = []

            for shape in slide.shapes:
                value = extract_pptx_shape_text(shape)
                if value:
                    texts.append(value)

            try:
                notes = slide.notes_slide.notes_text_frame.text
                notes = normalize_special_symbols((notes or "").strip())
                if notes:
                    texts.append("[备注]\n" + notes)
            except Exception:
                pass

            if texts:
                slides.append(f"[第 {index} 页]\n" + "\n".join(texts))

        return "\n\n".join(slides)
    except Exception:
        return ""


def read_docx_text(file_path: Path):
    """读取 DOCX 文本，包含段落和表格。需要安装 python-docx。"""
    try:
        from docx import Document
    except Exception:
        return ""

    try:
        document = Document(str(file_path))
        parts = []

        for paragraph in document.paragraphs:
            value = normalize_special_symbols((paragraph.text or "").strip())
            if value:
                parts.append(value)

        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    value = normalize_special_symbols((cell.text or "").strip())
                    if value:
                        cells.append(value)
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[表格 {table_index}]\n" + "\n".join(rows))

        return "\n".join(parts)
    except Exception:
        return ""


def normalize_ai_text(text: str, max_length: int = 1600):
    """文本清洗：保留页码线索，统一常见特殊符号，去掉过多空白。"""
    if not text:
        return ""

    cleaned = normalize_special_symbols(str(text))
    cleaned = (
        cleaned
        .replace("\r", "\n")
        .replace("\t", " ")
        .replace("\u3000", " ")
    )

    lines = []
    previous_blank = False

    for raw_line in cleaned.split("\n"):
        line = raw_line.strip()
        if not line:
            if not previous_blank:
                lines.append("")
            previous_blank = True
            continue

        # 去掉明显重复的超短空白，但保留页码 / 章节标记
        lines.append(line)
        previous_blank = False

    cleaned = "\n".join(lines).strip()
    return cleaned[:max_length]


def try_read_resource_text(resource: Optional[Resource]):
    """尽量读取用户上传的资料正文：txt / md / pdf / pptx / docx。"""
    if not resource:
        return ""

    file_path = Path(resource.file_path)
    if not file_path.exists() or not file_path.is_file():
        return ""

    suffix = file_path.suffix.lower()

    if suffix in {".txt", ".md"}:
        return read_plain_text_file(file_path)

    if suffix == ".pdf":
        return read_pdf_text(file_path)

    if suffix == ".pptx":
        return read_pptx_text(file_path)

    if suffix == ".docx":
        return read_docx_text(file_path)

    return ""


def infer_ai_keywords(course_name: str, resource_name: str, raw_text: str = ""):
    base = f"{course_name} {resource_name} {raw_text}".lower()
    keyword_map = [
        ("法理", ["法律关系", "权利义务", "法的渊源", "法治原则", "案例分析"]),
        ("法律", ["法律规范", "权利义务", "法律责任", "司法适用", "价值判断"]),
        ("python", ["数据结构", "函数封装", "异常处理", "文件读写", "算法思维"]),
        ("java", ["类与对象", "封装继承多态", "接口", "集合框架", "异常处理"]),
        ("离散", ["集合关系", "命题逻辑", "图论", "递推关系", "证明方法"]),
        ("宏观", ["总需求", "总供给", "GDP", "通货膨胀", "财政政策"]),
        ("经济", ["市场机制", "机会成本", "均衡分析", "政策影响", "数据解释"]),
        ("医学", ["核心概念", "病因机制", "诊断思路", "治疗原则", "复习要点"]),
        ("中医", ["阴阳五行", "脏腑经络", "辨证论治", "病因病机", "方药思路"]),
    ]

    for key, keywords in keyword_map:
        if key in base:
            return keywords

    return ["核心概念", "知识框架", "重点定义", "方法步骤", "复习问题"]


def infer_ai_sections(course_name: str, resource_name: str, raw_text: str = ""):
    keywords = infer_ai_keywords(course_name, resource_name, raw_text)
    text_hint = normalize_ai_text(raw_text, 600)

    material_summary = (
        "该资料主要围绕课程中的基础概念、关键方法和应用场景展开。"
        if not text_hint
        else "根据资料文本，内容可整理为若干概念、论证过程与复习问题。"
    )

    return {
        "keywords": keywords,
        "material_summary": material_summary,
        "text_hint": text_hint,
    }


def build_ai_note_content(course_name: str, resource_name: str, raw_text: str = "", note_style: str = "复习型"):
    sections = infer_ai_sections(course_name, resource_name, raw_text)
    keywords = sections["keywords"]

    keyword_lines = "\n".join([f"- **{item}**：建议结合教材定义、课堂例子和作业题进行复习。" for item in keywords])

    if sections["text_hint"]:
        source_excerpt = f"""
## 零、资料摘录参考

> 以下内容来自可读取资料文本，AI 已据此进行结构化整理：

{sections["text_hint"]}
"""
    else:
        source_excerpt = """
## 零、资料读取说明

当前资料可能是 PDF、PPT、图片或录音等格式，演示版先根据课程名、资料名和学习场景生成结构化笔记。
后续接入真实大模型与文档解析后，可直接读取资料正文、图片文字和音频转写内容。
"""

    return f"""# {course_name} · AI 结构化笔记

> 来源资料：{resource_name}
> 生成方式：NoteWhale AI powered 学习整理
> 笔记类型：{note_style}

{source_excerpt}

## 一、资料核心概览

{sections["material_summary"]}

本份笔记将资料拆分为“概念理解 → 重点归纳 → 复习任务 → 自测问题”四个层次，方便后续回看和考前复习。

## 二、关键词与核心概念

{keyword_lines}

## 三、知识框架

1. **概念层**
   - 先明确本节涉及的基本概念、定义和适用范围。
   - 对容易混淆的术语做对比整理。

2. **逻辑层**
   - 梳理概念之间的因果关系、推导关系或分类关系。
   - 将课堂内容转化为可复述的知识链条。

3. **应用层**
   - 结合例题、案例、作业或论文要求理解知识点。
   - 重点标记容易出现在考试、报告或课堂展示中的内容。

## 四、重点总结

- 本资料的复习重点不是单纯记忆，而是要能说明“为什么”和“如何使用”。
- 对于定义类内容，建议整理为：**定义 → 特征 → 例子 → 易错点**。
- 对于方法类内容，建议整理为：**步骤 → 条件 → 输出结果 → 常见错误**。
- 对于论文或报告类任务，建议提前整理参考文献、论证结构与格式要求。

## 五、复习建议

1. **第一遍：快速浏览**
   - 先看标题、目录、关键词，建立整体印象。

2. **第二遍：主动整理**
   - 用自己的话复述每个关键词。
   - 将资料内容补充到课堂笔记中。

3. **第三遍：自测检查**
   - 不看资料回答下方自测问题。
   - 对答不出的部分重新回到原资料定位。

## 六、自测问题

1. 这份资料最核心的 3 个概念是什么？
2. 这些概念之间有什么联系？
3. 哪些内容最可能出现在作业、论文或考试中？
4. 如果要向同学讲解这部分内容，应该按照什么顺序讲？
5. 当前还有哪些疑问需要回到教材或课堂中确认？

## 七、待补充

- [ ] 补充老师课堂强调内容
- [ ] 补充教材页码或参考文献
- [ ] 补充例题 / 案例
- [ ] 标记考试或作业重点
"""




def strip_markdown_fence(value: str):
    if not value:
        return ""

    text_value = str(value).strip()

    if text_value.startswith("```"):
        text_value = text_value.strip("`").strip()
        if text_value.lower().startswith("markdown"):
            text_value = text_value[8:].strip()
        elif text_value.lower().startswith("md"):
            text_value = text_value[2:].strip()

    return text_value.strip()


def post_process_ai_note(value: str):
    """模型输出后处理：去代码块、统一符号、修正部分 Markdown 空格。"""
    text_value = normalize_special_symbols(strip_markdown_fence(value))

    # 避免模型把标题写成 ####标题，编辑器不易识别
    import re
    text_value = re.sub(r"^(#{1,4})([^#\s])", r"\1 \2", text_value, flags=re.MULTILINE)

    # DeepSeek/Qwen 偶尔输出多余空行，保留段落但不要太散
    text_value = re.sub(r"\n{4,}", "\n\n\n", text_value)

    return text_value.strip()


def call_text_agent_for_course_note(
    course_name: str,
    resource_name: str,
    resource_text: str,
    note_style: str = "复习型",
):
    """
    调用真正的文本大模型生成课程资料笔记。
    需要配置：
    NOTEWHALE_TEXT_API_URL
    NOTEWHALE_TEXT_API_KEY
    NOTEWHALE_TEXT_MODEL

    默认接口按 OpenAI-compatible chat/completions 格式发送。
    """
    if not TEXT_API_URL or not TEXT_API_KEY or not TEXT_MODEL:
        raise HTTPException(
            status_code=503,
            detail="未配置文本模型智能体。请设置 NOTEWHALE_TEXT_API_URL / NOTEWHALE_TEXT_API_KEY / NOTEWHALE_TEXT_MODEL。",
        )

    clean_text = normalize_ai_text(resource_text, max_length=20000)

    if len(clean_text) < 80:
        raise HTTPException(
            status_code=422,
            detail="没有读取到足够的资料正文。请确认已安装 pypdf / python-pptx / python-docx，或上传可复制文字的资料。",
        )

    prompt = f"""
你是 NoteWhale 的“课程资料精读与复习笔记智能体”。

你的任务不是写大纲，而是把用户上传的课程资料加工成一份“可以直接复习、可以继续编辑、可以导出 PDF 的详细学习笔记”。
你需要尽量还原资料中的知识点、章节关系和页序线索，并把零散内容转化为可理解的复习材料。

课程：{course_name}
资料：{resource_name}
笔记类型：{note_style}

重要原则：
- 不要写空泛套话，例如“本章主要介绍相关知识”。
- 不要只列标题，不要只写三四条总结。
- 不能编造资料中没有的信息；资料中确实没有的内容写“资料中未给出”。
- 如果资料文本重复、断裂或像 PPT 短句，请主动补全逻辑衔接，但不要虚构新知识。
- 特殊符号、公式、逻辑符号优先输出 Unicode，可直接显示：≤ ≥ ≠ ≈ → ⇒ ⇔ ∑ ∫ √ ∈ ∉ ∪ ∩ ∀ ∃ ∴ ∵。
- 尽量避免复杂 LaTeX；如果必须写公式，用“公式 + 文字解释”的方式。

请按以下结构输出 Markdown：

# {course_name} · {resource_name} 详细复习笔记

## 1. 资料速读
写 5-8 条，说明资料类型、主题、范围、页码/章节线索、适合怎么复习。

## 2. 知识地图
用层级列表整理资料的知识结构。不是目录复刻，而是说明“概念之间的关系”。

## 3. 逐章 / 逐页精读
这是最重要的部分。
请按资料页码、章节或主题分成多个小节。每个小节都要包含：
- **本节讲什么**：用完整句子解释，不要只写关键词。
- **关键概念**：列出概念并解释含义。
- **逻辑关系**：说明概念之间如何连接，例如因果、分类、条件、推导、对比。
- **课堂可考点**：指出可能用于选择题、简答题、论述题或课堂讨论的点。
- **复习提醒**：指出容易忽略、容易混淆或需要回看原资料的地方。

## 4. 重点概念详解
选择资料中最重要的 8-15 个概念。
每个概念使用这个格式：
### 概念名
- **含义**：
- **为什么重要**：
- **与其他概念的关系**：
- **容易混淆点**：
- **一句话记忆**：

## 5. 易混淆点对比
用表格整理至少 4 组容易混淆的概念。如果资料中不足 4 组，请按资料实际内容输出。

表格列：
| 概念 A | 概念 B | 主要区别 | 判断方法 | 复习提示 |

## 6. 可回查索引
按页码、章节或资料线索写索引，方便用户回到原 PPT/PDF 查找。
格式示例：
- 第 1-3 页：……
- 第 4-8 页：……

## 7. 自测题
生成 10-15 道复习题，包含：
- 选择 / 判断倾向的问题
- 简答题
- 论述题
- 对比分析题

每道题后给出“答题要点”，不要只给题目。

## 8. 待补充清单
列出资料中不完整、需要老师课堂补充或需要教材查证的内容。

篇幅要求：
- 如果资料正文足够多，请输出 3000-5000 字。
- 如果资料正文较少，也要充分展开已有信息，但不要编造。
- “逐章 / 逐页精读”和“重点概念详解”必须是正文主体，不能只有大纲。

资料正文如下：
{clean_text}
"""

    body = {
        "model": TEXT_MODEL,
        "messages": [
            {
                "role": "system",
                "content": "你是一个严谨的课程资料精读、知识讲解与复习笔记生成智能体。你擅长把 PPT/PDF/Word 资料整理成详细、可复习、可编辑的 Markdown 笔记，而不是简单大纲。",
            },
            {
                "role": "user",
                "content": prompt,
            },
        ],
        "temperature": 0.35,
        "max_tokens": 6000,
    }

    request = urllib.request.Request(
        TEXT_API_URL,
        data=json.dumps(body).encode("utf-8"),
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {TEXT_API_KEY}",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(request, timeout=90) as response:
            result = json.loads(response.read().decode("utf-8"))
    except Exception as error:
        error_text = str(error)

        if "429" in error_text or "Too Many Requests" in error_text:
            raise HTTPException(
                status_code=429,
                detail=(
                    "文本模型请求过于频繁或资料过长，已触发免费额度限流。"
                    "请等待 1-3 分钟后重试，或换用页数更少的资料。"
                ),
            )

        raise HTTPException(status_code=502, detail=f"文本模型调用失败：{error}")

    try:
        content = result["choices"][0]["message"]["content"]
    except Exception:
        raise HTTPException(status_code=502, detail="文本模型返回格式异常，未找到 message.content。")

    return post_process_ai_note(content)



@app.post("/api/notes/generate")
def generate_note(
    payload: GenerateNoteRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """真正的 AI powered：读取资料正文后调用文本模型生成可编辑 Markdown 笔记。"""
    ensure_course_access(db, payload.courseId, current_user)

    resource = None

    if payload.resourceId:
        resource = (
            db.query(Resource)
            .filter(Resource.id == payload.resourceId, Resource.user_id == current_user.id)
            .first()
        )

    if not resource and payload.resourceName:
        resource = (
            db.query(Resource)
            .filter(
                Resource.user_id == current_user.id,
                Resource.course_id == payload.courseId,
                Resource.name == payload.resourceName,
            )
            .order_by(Resource.created_at.desc())
            .first()
        )

    file_resource_text = try_read_resource_text(resource)
    raw_resource_text = normalize_ai_text(payload.rawText, max_length=24000)
    resource_text = file_resource_text or raw_resource_text
    resource_name = resource.name if resource else payload.resourceName
    course_name = payload.courseName or "课程"

    content = call_text_agent_for_course_note(
        course_name=course_name,
        resource_name=resource_name or "课程资料",
        resource_text=resource_text,
        note_style=payload.noteStyle or "复习型",
    )

    now = datetime.utcnow()
    title = f"{course_name} · AI资料笔记"

    return {
        "id": int(now.timestamp() * 1000),
        "title": title,
        "content": content,
        "courseId": payload.courseId,
        "courseName": course_name,
        "source": resource_name or "课程资料",
        "aiGenerated": True,
        "createdAt": int(now.timestamp() * 1000),
        "updatedAt": int(now.timestamp() * 1000),
        "aiMeta": {
            "engine": TEXT_MODEL,
            "mode": "text-agent",
            "usedResourceText": bool(resource_text),
            "resourceId": resource.id if resource else None,
            "noteStyle": payload.noteStyle or "复习型",
        },
    }


def extract_json_object(text_value: str):
    """从视觉模型回复中提取 JSON 对象。"""
    if not text_value:
        raise ValueError("empty model response")

    text_value = text_value.strip()

    if text_value.startswith("```"):
        text_value = text_value.strip("`").strip()
        if text_value.lower().startswith("json"):
            text_value = text_value[4:].strip()

    start = text_value.find("{")
    end = text_value.rfind("}")

    if start == -1 or end == -1 or end <= start:
        raise ValueError("no json object")

    return json.loads(text_value[start : end + 1])


def normalize_agent_date(date_text: str):
    """统一模型返回的时间格式。"""
    if not date_text:
        return ""

    value = str(date_text).strip()
    value = value.replace("/", "-").replace("年", "-").replace("月", "-").replace("日", "")
    value = value.replace("T", " ")
    value = value.replace("：", ":")

    # 兼容 2026-7-10 18:00
    import re

    match = re.search(r"(\d{4})-(\d{1,2})-(\d{1,2})\s+(\d{1,2}):(\d{1,2})", value)
    if match:
        year, month, day, hour, minute = match.groups()
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d} {int(hour):02d}:{int(minute):02d}"

    match = re.search(r"(\d{1,2})-(\d{1,2})\s+(\d{1,2}):(\d{1,2})", value)
    if match:
        month, day, hour, minute = match.groups()
        year = datetime.utcnow().year
        return f"{year:04d}-{int(month):02d}-{int(day):02d} {int(hour):02d}:{int(minute):02d}"

    return value


def call_vision_agent_for_ddl(content: bytes, filename: str, course_name: str):
    """
    调用真正的视觉模型智能体识别 DDL 截图。
    需要配置：
    NOTEWHALE_VISION_API_URL
    NOTEWHALE_VISION_API_KEY
    NOTEWHALE_VISION_MODEL

    该函数按 OpenAI-compatible chat/completions 图片输入格式发送请求。
    """
    if not VISION_API_URL or not VISION_API_KEY or not VISION_MODEL:
        raise HTTPException(
            status_code=503,
            detail="未配置视觉模型智能体。请设置 NOTEWHALE_VISION_API_URL / NOTEWHALE_VISION_API_KEY / NOTEWHALE_VISION_MODEL。",
        )

    if not content:
        raise HTTPException(status_code=400, detail="图片内容为空")

    mime_type = mimetypes.guess_type(filename or "")[0] or "image/jpeg"
    image_base64 = base64.b64encode(content).decode("utf-8")
    image_url = f"data:{mime_type};base64,{image_base64}"

    current_year = datetime.utcnow().year

    prompt = f"""
你是 NoteWhale 的 DDL 截图识别智能体。

任务：只从截图中抽取“作业 / 论文 / 考试 / 提交任务”的结构化信息，不要总结，不要编造，不要扩写。

当前课程：{course_name or "未归属课程"}
当前年份：{current_year}

请严格输出一个 JSON 对象，不要 Markdown，不要解释。

字段要求：
{{
  "title": "任务标题。优先使用截图中出现的作业名、论文名、邮件标题或课程作业要求。不要写成笼统的课程名。",
  "date": "截止时间，格式必须是 YYYY-MM-DD HH:mm。如果截图只写月日，请用当前年份。",
  "platform": "提交平台或地点，例如：邮件提交、在线提交、教学平台、线下提交。",
  "note": "重要要求，保留原图关键信息，例如字数、格式、邮箱、文件名要求。不要写‘由AI生成’。",
  "confidence": 0.0到1.0之间的小数
}}

识别规则：
1. 如果图中有“7月10日18:00前提交”，date 应为 {current_year}-07-10 18:00。
2. 如果图中有邮箱或“发至”，platform 应优先为“邮件提交”。
3. 如果图中有“禁止使用AI”，note 中必须保留。
4. 如果图中有“邮件标题为”，title 或 note 中必须保留相关要求。
5. 如果确实看不清，不要乱编，字段可写“待确认”，confidence 降低。
"""

    body = {
        "model": VISION_MODEL,
        "messages": [
            {
                "role": "system",
                "content": "你是一个只输出 JSON 的课程 DDL 截图识别智能体。",
            },
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": image_url}},
                ],
            },
        ],
        "temperature": 0.1,
    }

    request = urllib.request.Request(
        VISION_API_URL,
        data=json.dumps(body).encode("utf-8"),
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {VISION_API_KEY}",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(request, timeout=45) as response:
            result = json.loads(response.read().decode("utf-8"))
    except Exception as error:
        raise HTTPException(status_code=502, detail=f"视觉模型调用失败：{error}")

    try:
        content_text = result["choices"][0]["message"]["content"]
        parsed = extract_json_object(content_text)
    except Exception:
        raise HTTPException(status_code=502, detail="视觉模型返回格式不是有效 JSON，请检查模型能力或提示词。")

    return {
        "title": str(parsed.get("title") or "待确认任务").strip(),
        "date": normalize_agent_date(parsed.get("date") or ""),
        "platform": str(parsed.get("platform") or "待确认").strip(),
        "note": str(parsed.get("note") or "").strip(),
        "confidence": float(parsed.get("confidence") or 0.0),
    }



@app.post("/api/ddls/recognize")
def recognize_ddl(
    payload: RecognizeDDLRequest,
    current_user: User = Depends(get_current_user),
):
    now = datetime.utcnow()

    return {
        "id": int(now.timestamp() * 1000),
        "title": f"{payload.courseName} 作业提交",
        "date": "2026-06-20 23:59",
        "courseId": payload.courseId,
        "courseName": payload.courseName,
        "platform": "教学平台",
        "note": "由截图识别生成，可继续手动修改。",
        "completed": False,
        "source": "图片识别",
        "createdAt": int(now.timestamp() * 1000),
    }


@app.post("/api/ddls/recognize-agent")
async def recognize_ddl_agent(
    file: UploadFile = File(...),
    courseId: Optional[int] = Query(None),
    courseName: str = Query("未归属课程"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """
    真正的视觉模型智能体识图接口：
    前端上传截图 -> 后端调用视觉模型 -> 返回结构化 DDL 草稿。
    """
    ensure_course_access(db, courseId, current_user)

    raw_name = Path(file.filename or "ddl_image.png").name
    content = await file.read()

    agent_result = call_vision_agent_for_ddl(
        content=content,
        filename=raw_name,
        course_name=courseName or "未归属课程",
    )

    now = datetime.utcnow()

    return {
        "id": int(now.timestamp() * 1000),
        "title": agent_result["title"],
        "date": agent_result["date"],
        "courseId": courseId,
        "courseName": courseName or "未归属课程",
        "platform": agent_result["platform"],
        "note": agent_result["note"],
        "completed": False,
        "source": "视觉模型智能体",
        "confidence": agent_result["confidence"],
        "imageName": raw_name,
        "createdAt": int(now.timestamp() * 1000),
        "aiMeta": {
            "engine": VISION_MODEL,
            "mode": "vision-agent",
        },
    }



@app.get("/api/resources")
def list_resources(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    resources = (
        db.query(Resource)
        .filter(Resource.user_id == current_user.id)
        .order_by(Resource.created_at.desc())
        .all()
    )
    return [resource_to_dict(resource) for resource in resources]


@app.post("/api/resources/upload")
async def upload_resource(
    file: UploadFile = File(...),
    courseId: Optional[int] = Query(None),
    courseName: str = Query(""),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    ensure_course_access(db, courseId, current_user)

    raw_name = Path(file.filename or "uploaded_file").name
    stored_name = f"{current_user.id}_{int(datetime.utcnow().timestamp() * 1000)}_{raw_name}"
    file_path = UPLOAD_DIR / stored_name

    content = await file.read()
    file_path.write_bytes(content)

    resource = Resource(
        name=raw_name,
        file_type=get_file_type(raw_name),
        file_path=str(file_path),
        size=len(content),
        course_id=courseId,
        course_name=courseName or "",
        user_id=current_user.id,
    )

    db.add(resource)
    db.commit()
    db.refresh(resource)

    return resource_to_dict(resource)


@app.delete("/api/resources/{resource_id}")
def delete_resource(
    resource_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    resource = (
        db.query(Resource)
        .filter(Resource.id == resource_id, Resource.user_id == current_user.id)
        .first()
    )
    if not resource:
        raise HTTPException(status_code=404, detail="资料不存在")

    file_path = Path(resource.file_path)

    db.delete(resource)
    db.commit()

    try:
        if file_path.exists() and file_path.is_file():
            file_path.unlink()
    except Exception:
        # 数据库删除已经完成，文件删除失败不影响接口返回。
        pass

    return {"ok": True, "deletedResourceId": resource_id}
